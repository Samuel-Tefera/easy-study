import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO
from sqlalchemy.orm import Session
from uuid import UUID
from app.repositories.document_repository import DocumentRepository
from app.services.llm_service import LLMService


class DocumentService:

    @staticmethod
    def extract_text(file_bytes: bytes, filename: str) -> str:
        suffix = filename.lower().split(".")[-1]

        if suffix == "pdf":
            return DocumentService._extract_pdf_text(file_bytes)
        elif suffix in ["pptx", "ppt"]:
            return DocumentService._extract_ppt_text(file_bytes)
        else:
            raise ValueError("Unsupported file format")

    @staticmethod
    def _extract_pdf_text(file_bytes: bytes) -> str:
        text = []
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for page in doc:
                text.append(page.get_text())
        return DocumentService._clean_text("\n".join(text))

    @staticmethod
    def _extract_ppt_text(file_bytes: bytes) -> str:
        prs = Presentation(BytesIO(file_bytes))
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            text.append(" | ".join(row_text))

        return DocumentService._clean_text("\n".join(text))

    @staticmethod
    def _clean_text(text: str) -> str:
        # Remove NUL characters (PostgreSQL doesn't support them in string literals)
        text = text.replace("\0", "")
        text = text.replace("\x00", "")
        text = text.replace("\n\n", "\n")
        return text.strip()

    @staticmethod
    def generate_document_summary(db: Session, document_id: UUID) -> str:
        # Check if already generated
        document = DocumentRepository.get_document_by_id(db, document_id)
        if not document:
            raise ValueError(f"Document {document_id} not found")
        if document.summary:
            return document.summary

        chunks = DocumentRepository.get_document_chunks(db, document_id)
        if not chunks:
            raise ValueError(f"No text extracted for document {document_id}")

        # Combine all parts
        combined_text = "\n\n".join([chunk.text for chunk in chunks])

        # Generate summary using LLM service
        summary_markdown = LLMService.generate_summary(combined_text)

        # Save to DB
        DocumentRepository.update_document_summary(db, document_id, summary_markdown)

        return summary_markdown

